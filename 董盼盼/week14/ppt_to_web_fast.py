import argparse
import html
import re
import sys
from pathlib import Path

TEMPLATE = """<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>__TITLE__</title>
<style>
:root{--bg:#0f172a;--sf:#1e293b;--sh:#334155;--pr:#6366f1;--pl:#818cf8;--pd:#4f46e5;--ac:#06b6d4;--al:#22d3ee;--gr:#10b981;--gl:#34d399;--or:#f59e0b;--rd:#ef4444;--pk:#ec4899;--tx:#f1f5f9;--tm:#94a3b8;--td:#64748b;--bd:#334155;--cb:#1a2332;--sd:0 4px 24px rgba(0,0,0,.3)}
*{margin:0;padding:0;box-sizing:border-box}
body{font-family:'Segoe UI',system-ui,sans-serif;background:var(--bg);color:var(--tx);line-height:1.6;overflow-x:hidden}
.nav{position:fixed;top:0;left:0;right:0;height:56px;background:rgba(15,23,42,.95);backdrop-filter:blur(10px);border-bottom:1px solid var(--bd);z-index:100;display:flex;align-items:center;padding:0 24px;gap:12px}
.brand{font-size:18px;font-weight:700;background:linear-gradient(135deg,var(--pl),var(--al));-webkit-background-clip:text;-webkit-text-fill-color:transparent;white-space:nowrap}
.links{display:flex;gap:4px;margin-left:24px;overflow-x:auto}
.links a{color:var(--tm);text-decoration:none;font-size:13px;padding:6px 12px;border-radius:6px;white-space:nowrap;transition:all .2s}
.links a:hover,.links a.on{color:var(--tx);background:var(--sh)}
.side{position:fixed;top:56px;left:0;bottom:0;width:220px;background:var(--sf);border-right:1px solid var(--bd);overflow-y:auto;padding:16px 0;z-index:50}
.si{display:flex;align-items:center;gap:8px;padding:10px 20px;cursor:pointer;color:var(--tm);font-size:13px;transition:all .2s;border-left:3px solid transparent}
.si:hover{color:var(--tx);background:var(--sh)}
.si.on{color:var(--pl);border-left-color:var(--pr);background:rgba(99,102,241,.08)}
.si .n{font-size:11px;color:var(--td);min-width:24px}
.main{margin-left:220px;margin-top:56px;padding:40px 48px;max-width:1000px}
.cover{text-align:center;padding:60px 0 80px}
.cover .tags{display:flex;justify-content:center;gap:12px;margin-bottom:32px;flex-wrap:wrap}
.cover .tag{padding:6px 16px;border-radius:20px;font-size:13px;font-weight:600}
.cover .tag.purple{background:rgba(99,102,241,.15);color:var(--pl)}
.cover .tag.cyan{background:rgba(6,182,212,.15);color:var(--al)}
.cover .tag.green{background:rgba(16,185,129,.15);color:var(--gl)}
.cover .tag.pink{background:rgba(236,72,153,.15);color:var(--pk)}
.cover h1{font-size:48px;font-weight:800;margin-bottom:16px;background:linear-gradient(135deg,var(--pl),var(--al),var(--gl));-webkit-background-clip:text;-webkit-text-fill-color:transparent}
.cover .sub{font-size:22px;color:var(--tm);margin-bottom:12px}
.cover .desc{font-size:16px;color:var(--td)}
.sh{margin-bottom:32px;padding-bottom:16px;border-bottom:1px solid var(--bd)}
.sh .pn{font-size:14px;font-weight:600;color:var(--pl);text-transform:uppercase;letter-spacing:2px;margin-bottom:8px}
.sh h2{font-size:32px;font-weight:700}
.sh .s{font-size:15px;color:var(--tm);margin-top:8px}
.card{background:var(--sf);border-radius:12px;padding:28px;margin-bottom:24px;border:1px solid var(--bd);transition:border-color .2s}
.card:hover{border-color:var(--pr)}
.card h3{font-size:20px;font-weight:600;margin-bottom:16px}
.card h4{font-size:15px;font-weight:700;margin:20px 0 10px;color:var(--al);padding-left:10px;border-left:3px solid var(--ac)}
.card p{font-size:14px;color:var(--tm);margin-bottom:8px}
.code{background:var(--cb);border-radius:8px;padding:20px;font-family:'Cascadia Code','Fira Code',Consolas,monospace;font-size:13px;line-height:1.7;overflow-x:auto;border:1px solid var(--bd);color:#e2e8f0;white-space:pre-wrap;word-break:break-word;margin:12px 0}
.tbl{width:100%;border-collapse:collapse;margin:12px 0}
.tbl th{text-align:left;padding:12px 16px;font-size:14px;color:var(--pl);border-bottom:2px solid var(--bd)}
.tbl td{padding:12px 16px;font-size:14px;color:var(--tm);border-bottom:1px solid var(--bd)}
.tbl tr:hover td{background:rgba(99,102,241,.05)}
.hb{background:rgba(99,102,241,.1);border-left:4px solid var(--pr);padding:16px 20px;border-radius:0 8px 8px 0;margin:16px 0;font-size:15px;color:var(--tx)}
.hb.g{background:rgba(16,185,129,.1);border-left-color:var(--gr)}
.hb.o{background:rgba(245,158,11,.1);border-left-color:var(--or)}
.grid{display:grid;gap:16px;margin:16px 0}
.grid.c2{grid-template-columns:repeat(2,1fr)}
.grid.c3{grid-template-columns:repeat(3,1fr)}
.grid.c4{grid-template-columns:repeat(4,1fr)}
.gi{background:rgba(15,23,42,.5);border-radius:10px;padding:20px;border:1px solid var(--bd);transition:all .2s}
.gi:hover{border-color:var(--pr);transform:translateY(-2px);box-shadow:var(--sd)}
.gi .gn{display:inline-block;font-size:12px;font-weight:700;color:var(--pl);background:rgba(99,102,241,.15);padding:2px 10px;border-radius:10px;margin-bottom:10px}
.gi .gt{font-size:16px;font-weight:700;margin-bottom:6px;color:var(--tx)}
.gi .gl{font-size:13px;color:var(--al);margin-bottom:8px;font-weight:600}
.gi .gd{font-size:13px;color:var(--tm);line-height:1.5}
.arch{display:flex;flex-direction:column;align-items:center;gap:0;margin:16px 0;padding:20px;background:rgba(15,23,42,.5);border-radius:12px;border:1px solid var(--bd)}
.al{display:flex;align-items:center;gap:12px;padding:12px 24px;background:var(--sh);border-radius:8px;border:1px solid var(--bd);min-width:280px;justify-content:center}
.al .an{font-size:12px;font-weight:700;color:var(--pl);background:rgba(99,102,241,.15);padding:2px 10px;border-radius:10px}
.al .at{font-size:15px;font-weight:600;color:var(--tx)}
.aa{color:var(--al);font-size:20px;margin:4px 0}
.sl{display:flex;flex-direction:column;gap:8px;margin:12px 0}
.sli{display:flex;align-items:center;gap:16px;padding:12px 16px;background:rgba(15,23,42,.5);border-radius:8px;border:1px solid var(--bd)}
.sli .sv{font-size:22px;font-weight:800;min-width:80px;background:linear-gradient(135deg,var(--pl),var(--al));-webkit-background-clip:text;-webkit-text-fill-color:transparent}
.sli .sl-l{font-size:14px;font-weight:600;color:var(--tx)}
.sli .sl-d{font-size:13px;color:var(--tm);margin-left:8px}
.flow{display:flex;align-items:center;gap:8px;flex-wrap:wrap;margin:12px 0}
.fs{padding:10px 20px;background:var(--sh);border-radius:8px;font-size:14px;color:var(--tx);border:1px solid var(--bd)}
.fa{color:var(--td);font-size:18px}
.ct{width:100%;border-collapse:collapse;margin:12px 0}
.ct th{text-align:left;padding:12px 16px;font-size:14px;color:var(--pl);border-bottom:2px solid var(--bd)}
.ct td{padding:12px 16px;font-size:14px;color:var(--tm);border-bottom:1px solid var(--bd)}
.ct td:first-child{color:var(--tx);font-weight:600;white-space:nowrap;width:30%}
.ct tr:hover td{background:rgba(99,102,241,.05)}
.sr{display:flex;gap:16px;flex-wrap:wrap;margin:12px 0}
.sr .st{flex:1;min-width:140px;background:var(--sf);border-radius:12px;padding:20px;text-align:center;border:1px solid var(--bd)}
.sr .st .sv{font-size:28px;font-weight:800;background:linear-gradient(135deg,var(--pl),var(--al));-webkit-background-clip:text;-webkit-text-fill-color:transparent}
.sr .st .sv.r{background:linear-gradient(135deg,var(--rd),var(--or));-webkit-background-clip:text;-webkit-text-fill-color:transparent}
.sr .st .sv.g{background:linear-gradient(135deg,var(--gl),var(--ac));-webkit-background-clip:text;-webkit-text-fill-color:transparent}
.sr .st .sl{font-size:13px;color:var(--tm);margin-top:6px}
.sr .st .sl b{color:var(--tx);font-weight:600;display:block;margin-bottom:2px}
.cg{display:grid;grid-template-columns:repeat(2,1fr);gap:16px;margin:12px 0}
.ci{background:var(--sf);border-radius:10px;padding:20px;border:1px solid var(--bd);border-left:4px solid var(--pr)}
.ci.h{border-left-color:var(--gr)}
.ci .ct-t{font-weight:700;font-size:15px;margin-bottom:12px;color:var(--pl)}
.ci.h .ct-t{color:var(--gl)}
.ci .ct-p{color:var(--gl);font-size:13px;margin-bottom:6px}
.ci .ct-c{color:var(--rd);font-size:13px;margin-bottom:6px}
.ci .ct-f{color:var(--tm);font-size:13px}
.vg{display:grid;grid-template-columns:repeat(3,1fr);gap:16px;margin:12px 0}
.vc{background:var(--sf);border-radius:12px;padding:24px;text-align:center;border:1px solid var(--bd);border-top:3px solid var(--pr)}
.vc .ic{font-size:28px;margin-bottom:8px}
.vc .bg{font-size:18px;font-weight:800;margin-bottom:6px}
.vc .bg.g{color:var(--gl)}
.vc .bg.p{color:var(--pl)}
.vc .bg.c{color:var(--al)}
.vc .tt{font-size:14px;font-weight:600;color:var(--tx);margin-bottom:8px}
.vc p{font-size:13px;color:var(--tm);line-height:1.6}
.chg{display:grid;grid-template-columns:repeat(2,1fr);gap:16px;margin:12px 0}
.chc{background:var(--sf);border-radius:10px;padding:20px;border:1px solid var(--bd);border-left:4px solid var(--or)}
.chc h4{font-size:15px;font-weight:700;color:var(--pl);margin-bottom:8px}
.chc p{font-size:13px;color:var(--tm);line-height:1.5}
.end{text-align:center;padding:60px 0}
.end h2{font-size:36px;font-weight:800;margin-bottom:16px;background:linear-gradient(135deg,var(--pl),var(--gl));-webkit-background-clip:text;-webkit-text-fill-color:transparent}
::-webkit-scrollbar{width:6px;height:6px}
::-webkit-scrollbar-track{background:var(--bg)}
::-webkit-scrollbar-thumb{background:var(--sh);border-radius:3px}
@media(max-width:768px){.side{display:none}.main{margin-left:0;padding:24px 16px}.grid.c2,.grid.c3,.grid.c4{grid-template-columns:1fr}}
</style>
</head>
<body>
<nav class="nav"><div class="brand">__BRAND__</div><div class="links">__LINKS__</div></nav>
<aside class="side">__SIDE__</aside>
<main class="main">__CONTENT__</main>
<script>
function go(id){const e=document.getElementById(id);if(e){const t=e.getBoundingClientRect().top+window.pageYOffset-70;window.scrollTo({top:t,behavior:'smooth'})}}
const sec=document.querySelectorAll('section[id]'),si=document.querySelectorAll('.si'),nl=document.querySelectorAll('.links a');
window.addEventListener('scroll',()=>{let cur='';sec.forEach(s=>{if(s.getBoundingClientRect().top<=120)cur=s.id});si.forEach((i,k)=>{const ids=[...sec].map(s=>s.id);i.classList.toggle('on',ids[k]===cur)});nl.forEach(l=>{l.classList.toggle('on',l.getAttribute('href')==='#'+cur)})});
</script>
</body>
</html>"""


COLORS = ["g", "p", "c"]


def esc(t):
    return html.escape(str(t))


def is_watermark(text):
    return "八斗学院" in text or "盗版" in text


# ===== 内容类型识别 =====

def is_code(text):
    indicators = [
        text.startswith("{"), text.startswith("---"), text.startswith("def "),
        text.startswith("import "), text.startswith("from "), text.startswith("class "),
        text.startswith("#!/"), text.startswith("```"),
        "function " in text[:20], "const " in text[:20], "return " in text[:20],
        "pip install" in text, "npm install" in text, "python -m" in text,
        "git " in text[:10],
        re.search(r"^\s*(if|for|while|try|with)\s", text) is not None,
    ]
    return any(indicators)


def is_highlight(text):
    keywords = ["核心矛盾", "核心问题", "核心思想", "核心价值", "核心机制", "本质特征",
                "关键点", "关键在于", "关键是", "总结", "综上", "因此", "本质上",
                "重要", "注意", "警告", "提示", "原则", "选型依据", "形态之间"]
    return any(kw in text[:25] for kw in keywords) and len(text) > 15


def is_stat_line(text):
    t = text.strip()
    # 纯数字（如 01、02）不算统计值，可能是步骤编号
    if re.match(r"^\d{1,2}$", t):
        return False
    return bool(re.match(r"^[~]?\d+[%x倍]?\+?$", t))


def is_combined_stat(text):
    t = text.strip()
    # 必须包含单位指示符（%、x、倍、+），排除 "01\n用户发消息" 这类步骤编号
    return bool(re.match(r"^[~]?\d+[%x倍+]\s+\S+", t))


def is_step_number(text):
    """检测步骤编号文本，如 '01\\n用户发消息'"""
    t = text.strip()
    return bool(re.match(r"^\d{1,2}[\s\n]+[^\d]", t))


def parse_combined_stat(text):
    t = text.strip()
    m = re.match(r"^([~]?\d+[%x倍]?\+?)\s+(.+)$", t)
    if m:
        return m.group(1), m.group(2)
    return t, ""


def is_stage_marker(text):
    t = text.strip()
    if re.match(r"^\d{1,2}$", t):
        return True
    if "\n" in t and re.match(r"^\d{1,2}\n", t):
        return True
    patterns = [r"^阶段[一二三四五六七八九十\d]+$", r"^[Pp]hase\s*\d+$",
                r"^第[一二三四五六七八九十\d]+[阶段步]$", r"^Step\s*\d+$"]
    return any(re.match(p, t) for p in patterns)


def is_year(text):
    return bool(re.match(r"^\d{4}(-\d{2,4})?$", text.strip()))


def is_short_title(text):
    t = text.strip()
    if len(t) < 2 or len(t) > 15:
        return False
    if any(c in t for c in "。，；：！？、\n"):
        return False
    return True


def is_description(text):
    t = text.strip()
    if len(t) < 8:
        return False
    return "\n" in t or "。" in t or "，" in t or "、" in t or len(t) > 12


SUBTITLE_KEYWORDS = [
    "核心价值", "核心机制", "核心思想", "本质特征", "架构层级", "流程示意",
    "方案对比", "带来的", "代价", "收益", "构成", "结构", "定义",
    "局限性", "关系", "量化", "生命周期", "变化", "形态", "调用形式",
    "单次交互", "定义方式", "与FC关系", "触发检测", "执行层", "触发层",
    "常驻层", "MEMORY", "Context", "结构说明", "设计原则", "核心价值主张",
    "挑战", "开放问题",
]


def is_subtitle(text):
    t = text.strip()
    if len(t) < 2 or len(t) > 25:
        return False
    if any(c in t for c in "。，；！？\n"):
        return False
    if is_stat_line(t) or is_code(t) or is_highlight(t) or is_stage_marker(t) or is_year(t):
        return False
    if t.endswith("：") or t.endswith(":"):
        return True
    if any(kw in t for kw in SUBTITLE_KEYWORDS):
        return True
    return False


def is_arch_layer(text):
    t = text.strip()
    if len(t) > 30:
        return False
    if any(c in t for c in "。，；！？"):
        return False
    # 排除子标题（如 "MCP 架构层级"、"RAG 流程示意"）
    if is_subtitle(t):
        return False
    arch_keywords = ["Client", "Server", "Agent", "API", "DB", "LLM", "MCP",
                     "层", "顶层", "底层", "上层", "下层", "前端", "后端",
                     "输入", "输出", "处理", "存储", "实际工具", "Function",
                     "Call", "RAG"]
    if any(kw in t for kw in arch_keywords):
        return True
    if "/" in t and len(t) <= 25:
        return True
    return False


def has_arrow(text):
    return "→" in text or "=>" in text


# ===== 渲染函数 =====

def render_table(data):
    if not data:
        return ""
    rows = []
    for i, row in enumerate(data):
        tag = "th" if i == 0 else "td"
        cells = "".join(f"<{tag}>{esc(c)}</{tag}>" for c in row)
        rows.append(f"<tr>{cells}</tr>")
    return f'<table class="tbl">{"".join(rows)}</table>'


def render_code(text):
    return f'<div class="code">{esc(text)}</div>'


def render_highlight(text):
    cls = "g" if any(kw in text for kw in ["收益", "优势", "节省", "提升"]) else ""
    if any(kw in text for kw in ["代价", "不足", "警告", "问题", "挑战"]):
        cls = "o"
    return f'<div class="hb {cls}">{esc(text)}</div>'


def render_subtitle(text):
    return f"<h4>{esc(text)}</h4>"


def render_card_grid(items):
    n = len(items)
    cols = 4 if n >= 4 else (3 if n == 3 else 2)
    items_html = []
    for item in items:
        parts = []
        if item.get("num"):
            parts.append(f'<span class="gn">{esc(item["num"])}</span>')
        if item.get("title"):
            parts.append(f'<div class="gt">{esc(item["title"])}</div>')
        if item.get("label"):
            parts.append(f'<div class="gl">{esc(item["label"])}</div>')
        if item.get("desc"):
            parts.append(f'<div class="gd">{esc(item["desc"])}</div>')
        if item.get("desc2"):
            parts.append(f'<div class="gd" style="margin-top:8px;color:var(--al)">{esc(item["desc2"])}</div>')
        items_html.append(f'<div class="gi">{"".join(parts)}</div>')
    return f'<div class="grid c{cols}">{"".join(items_html)}</div>'


def render_architecture(layers):
    items_html = []
    for i, layer in enumerate(layers):
        if i > 0:
            items_html.append('<div class="aa">↓</div>')
        items_html.append(
            f'<div class="al"><span class="an">L{i+1}</span>'
            f'<span class="at">{esc(layer)}</span></div>'
        )
    return f'<div class="arch">{"".join(items_html)}</div>'


def render_stat_list(stats):
    items_html = []
    for s in stats:
        desc_html = ""
        if s.get("desc"):
            desc_html = f'<div class="sl-d">{esc(s["desc"])}</div>'
        items_html.append(
            f'<div class="sli"><div class="sv">{esc(s["value"])}</div>'
            f'<div class="sl-l">{esc(s["label"])}</div>{desc_html}</div>'
        )
    return f'<div class="sl">{"".join(items_html)}</div>'


def render_flow(steps):
    items_html = []
    for i, step in enumerate(steps):
        if i > 0:
            items_html.append('<span class="fa">→</span>')
        items_html.append(f'<div class="fs">{esc(step)}</div>')
    return f'<div class="flow">{"".join(items_html)}</div>'


# ===== 模式匹配 =====

def match_card_item(texts, j):
    """匹配单个卡片项，返回 (consumed, item_dict)"""
    if j >= len(texts):
        return 0, None
    t1 = texts[j]
    if is_subtitle(t1):
        return 0, None
    # 架构层级交由 try_architecture 处理
    if is_arch_layer(t1):
        return 0, None

    # 模式 A: 编号 + 短标题 + 描述 + "典型场景" + 描述2（5字段）
    if is_stage_marker(t1) and j + 4 < len(texts):
        t2, t3, t4, t5 = texts[j+1], texts[j+2], texts[j+3], texts[j+4]
        if (is_short_title(t2) and is_description(t3)
            and t4 == "典型场景" and is_description(t5)):
            return 5, {"num": t1, "title": t2, "desc": t3, "label": "典型场景", "desc2": t5}

    # 模式 B: 阶段标识 + 短标题 + 年份 + 描述（4字段）
    if is_stage_marker(t1) and j + 3 < len(texts):
        t2, t3, t4 = texts[j+1], texts[j+2], texts[j+3]
        if (is_short_title(t2) and is_year(t3) and is_description(t4)
            and not is_code(t2) and not is_code(t4)):
            return 4, {"num": t1, "title": t2, "label": t3, "desc": t4}

    # 模式 C: 短标题 + 短标签 + 描述（3字段）
    if is_short_title(t1) and j + 2 < len(texts):
        t2, t3 = texts[j+1], texts[j+2]
        if (is_short_title(t2) and is_description(t3)
            and not is_subtitle(t2)
            and not is_code(t1) and not is_code(t2) and not is_code(t3)
            and not is_highlight(t3) and not is_stat_line(t1)):
            return 3, {"title": t1, "label": t2, "desc": t3}

    # 模式 D: 短标题 + 描述（2字段）
    if is_short_title(t1) and j + 1 < len(texts):
        t2 = texts[j+1]
        if (is_description(t2)
            and not is_code(t1) and not is_code(t2)
            and not is_highlight(t2) and not is_stat_line(t1)):
            return 2, {"title": t1, "desc": t2}

    # 模式 E: 编号+标题（带换行，如 "1. 单一职责"）
    if re.match(r"^\d+\.\s", t1) and j + 1 < len(texts):
        title = re.sub(r"^\d+\.\s*", "", t1)
        t2 = texts[j+1]
        if is_description(t2) and not is_code(t2) and not is_highlight(t2):
            num = re.match(r"^(\d+)\.", t1).group(1)
            return 2, {"num": num, "title": title, "desc": t2}

    return 0, None


def try_card_grid(texts, i):
    items = []
    j = i
    while j < len(texts):
        consumed, item = match_card_item(texts, j)
        if consumed > 0 and item:
            items.append(item)
            j += consumed
        else:
            break
    if len(items) >= 2:
        return j - i, render_card_grid(items)
    return 0, ""


def try_architecture(texts, i):
    layers = []
    j = i
    while j < len(texts):
        t = texts[j]
        if is_arch_layer(t):
            layers.append(t)
            j += 1
        else:
            break
    if len(layers) >= 3:
        return j - i, render_architecture(layers)
    return 0, ""


def try_stat_list(texts, i):
    stat_items = []
    j = i
    while j < len(texts):
        t = texts[j]

        # 模式 A: 组合统计行 "~15%  用户指令" + 描述
        if is_combined_stat(t):
            value, label = parse_combined_stat(t)
            if j + 1 < len(texts):
                t_next = texts[j+1]
                # 接受任何非统计、非代码、非子标题的文本作为描述
                if (not is_stat_line(t_next) and not is_combined_stat(t_next)
                    and not is_code(t_next) and not is_subtitle(t_next)
                    and not is_highlight(t_next) and len(t_next) <= 50):
                    stat_items.append({"value": value, "label": label, "desc": t_next})
                    j += 2
                    continue
            stat_items.append({"value": value, "label": label})
            j += 1
            continue

        # 模式 B: 纯统计 + 标签 + 描述
        if is_stat_line(t) and j + 2 < len(texts):
            t2, t3 = texts[j+1], texts[j+2]
            if (is_short_title(t2) and is_description(t3)
                and not is_stat_line(t2) and not is_code(t2)):
                stat_items.append({"value": t, "label": t2, "desc": t3})
                j += 3
                continue

        # 模式 C: 纯统计 + 标签
        if is_stat_line(t) and j + 1 < len(texts):
            t2 = texts[j+1]
            if (is_short_title(t2) and not is_stat_line(t2) and not is_code(t2)
                and not is_combined_stat(t2)):
                stat_items.append({"value": t, "label": t2})
                j += 2
                continue

        break

    if len(stat_items) >= 2:
        return j - i, render_stat_list(stat_items)
    return 0, ""


def try_table(texts, i):
    """匹配表格：表头行（3+ 短文本）+ 数据行"""
    headers = []
    j = i
    while j < len(texts) and is_short_title(texts[j]):
        headers.append(texts[j])
        j += 1

    if len(headers) < 3:
        return 0, ""

    rows = []
    current_row = []
    while j < len(texts):
        t = texts[j]
        if is_subtitle(t) or is_highlight(t):
            break
        if is_short_title(t) and len(current_row) >= len(headers) - 1:
            while len(current_row) < len(headers):
                current_row.append("")
            rows.append(current_row)
            current_row = [t]
            j += 1
            continue
        current_row.append(t)
        j += 1
        if len(current_row) == len(headers):
            rows.append(current_row)
            current_row = []

    if current_row:
        while len(current_row) < len(headers):
            current_row.append("")
        rows.append(current_row)

    if len(rows) < 2:
        return 0, ""

    table_data = [headers] + rows
    return j - i, render_table(table_data)


def try_flow_steps(texts, i):
    steps = []
    j = i
    while j < len(texts):
        t = texts[j]
        if (len(t) <= 20 and not is_stat_line(t) and not is_code(t)
            and not is_highlight(t) and not is_subtitle(t)
            and not is_description(t)):
            steps.append(t)
            j += 1
        else:
            break
    if len(steps) >= 3:
        return j - i, render_flow(steps)
    return 0, ""


def try_step_flow(texts, i):
    """匹配步骤流程：'01\\n标题' 格式的步骤序列"""
    steps = []
    j = i
    while j < len(texts):
        t = texts[j]
        if is_step_number(t):
            # 提取编号和标题
            m = re.match(r"^(\d{1,2})[\s\n]+(.+)$", t.strip())
            if m:
                steps.append(f"{m.group(1)} {m.group(2)}")
            else:
                steps.append(t)
            j += 1
        else:
            break
    if len(steps) >= 2:
        return j - i, render_flow(steps)
    return 0, ""


def try_concept_table(texts, i):
    """匹配概念表：2+ 对 (短标题 + 描述)"""
    pairs = []
    j = i
    while j + 1 < len(texts):
        t1, t2 = texts[j], texts[j+1]
        # 短标题 + 描述对
        # 注意：即使 t1 是子标题关键字（如"局限性"），只要后跟描述且已在表中，就作为键
        if (is_short_title(t1) and is_description(t2)
            and not is_arch_layer(t1)
            and not is_stat_line(t1) and not is_combined_stat(t1)
            and not is_code(t1) and not is_code(t2)
            and not is_highlight(t2) and not is_stage_marker(t1)
            and len(t1) <= 12
            and (pairs or not is_subtitle(t1))):
            pairs.append((t1, t2))
            j += 2
            continue
        break
    if len(pairs) >= 3:
        rows = "".join(f"<tr><td>{esc(k)}</td><td>{esc(v)}</td></tr>" for k, v in pairs)
        return j - i, f'<table class="ct"><tbody>{rows}</tbody></table>'
    return 0, ""


def try_stat_row(texts, i):
    """匹配水平统计行：2+ 个统计项（值+标签[+描述]）"""
    items = []
    j = i
    while j < len(texts):
        t = texts[j]
        # 跳过步骤编号文本（如 "01\n用户发消息"）
        if is_step_number(t):
            break
        # 模式 A: 组合统计 "~15%  用户指令" + 可选描述
        if is_combined_stat(t):
            value, label = parse_combined_stat(t)
            desc = ""
            if j + 1 < len(texts):
                t_next = texts[j+1]
                if (not is_stat_line(t_next) and not is_combined_stat(t_next)
                    and not is_code(t_next) and not is_subtitle(t_next)
                    and not is_highlight(t_next) and not is_step_number(t_next)
                    and len(t_next) <= 30):
                    desc = t_next
                    j += 2
                else:
                    j += 1
            else:
                j += 1
            items.append({"value": value, "label": label, "desc": desc})
            continue
        # 模式 B: 纯统计 + 标签
        if is_stat_line(t) and j + 1 < len(texts):
            t2 = texts[j+1]
            if (is_short_title(t2) and not is_stat_line(t2) and not is_code(t2)
                and not is_combined_stat(t2) and not is_subtitle(t2)):
                items.append({"value": t, "label": t2, "desc": ""})
                j += 2
                continue
        break
    if len(items) >= 3:
        cells = ""
        for it in items:
            desc_html = f"<br>{esc(it['desc'])}" if it.get("desc") else ""
            cells += f'<div class="st"><div class="sv">{esc(it["value"])}</div><div class="sl"><b>{esc(it["label"])}</b>{desc_html}</div></div>'
        return j - i, f'<div class="sr">{cells}</div>'
    return 0, ""


def try_compare_grid(texts, i):
    """匹配策略对比卡片：标题 + 优势 + 优势描述 + 不足 + 不足描述 + 适用 + 适用描述"""
    items = []
    j = i
    while j + 6 <= len(texts):
        t1 = texts[j]
        if not is_short_title(t1) or is_subtitle(t1) or is_arch_layer(t1):
            break
        # 必须是 "标题 + 优势 + 描述 + 不足 + 描述 + 适用 + 描述" 模式
        if (texts[j+1] == "优势" and texts[j+3] == "不足" and texts[j+5] == "适用"):
            items.append({
                "title": t1,
                "pro": texts[j+2],
                "con": texts[j+4],
                "fit": texts[j+6] if j + 6 < len(texts) else ""
            })
            j += 7
            continue
        break
    if len(items) >= 2:
        cells = ""
        for it in items:
            is_highlight_item = "渐进式披露" in it["title"]
            cls = "ci h" if is_highlight_item else "ci"
            cells += (
                f'<div class="{cls}">'
                f'<div class="ct-t">{esc(it["title"])}</div>'
                f'<div class="ct-p">✓ {esc(it.get("pro",""))}</div>'
                f'<div class="ct-c">✗ {esc(it.get("con",""))}</div>'
                f'<div class="ct-f">适用：{esc(it.get("fit",""))}</div>'
                f'</div>'
            )
        return j - i, f'<div class="cg">{cells}</div>'
    return 0, ""


def try_value_grid(texts, i):
    """匹配价值卡片：[副标题 + 大字 + 标签 + 多行描述] × 3
    每个 card: subtitle(2-12字) + big(2-10字，如"节省 60-90%"/"可复用"/"无上限") + label(2-10字) + 3行描述"""
    items = []
    j = i
    while j + 3 < len(texts):
        t1 = texts[j]
        # t1 必须是短标题（副标题，如"Context 效率"/"能力模块化"/"系统可扩展"）
        if not is_short_title(t1) or is_arch_layer(t1) or is_code(t1):
            break
        if len(t1) > 12:
            break
        # t2 是大字值（短，2-10字）
        t2 = texts[j+1]
        if not is_short_title(t2) or is_subtitle(t2) or len(t2) > 12:
            break
        # t3 是标签（短，2-10字）
        t3 = texts[j+2]
        if not is_short_title(t3) or is_subtitle(t3) or len(t3) > 12:
            break
        # 收集描述行（直到下一个副标题或模式结束）
        desc_lines = []
        k = j + 3
        while k < len(texts):
            t = texts[k]
            # 遇到下一个 card 的开始（短副标题≤8字 + 短标题 + 短标题 模式）
            # 副标题很短（如"能力模块化"/"系统可扩展"），描述行通常更长
            if (is_short_title(t) and not is_subtitle(t) and k + 2 < len(texts)
                and is_short_title(texts[k+1]) and not is_subtitle(texts[k+1])
                and is_short_title(texts[k+2]) and not is_subtitle(texts[k+2])
                and len(t) <= 8):
                break
            # 遇到明显的高亮总结或代码
            if is_highlight(t) or is_code(t):
                break
            # 遇到总结性句子（含"="或"从工具调用"等）停止
            if ("=" in t and len(t) > 15) or ("从工具调用" in t):
                break
            desc_lines.append(t)
            k += 1
            # 最多收集 4 行描述
            if len(desc_lines) >= 4:
                break
        if desc_lines:
            items.append({
                "title": t1,
                "big": t2,
                "label": t3,
                "desc": "<br>".join(esc(d) for d in desc_lines)
            })
            j = k
            continue
        break
    if len(items) >= 2:
        cells = ""
        for idx_it, it in enumerate(items):
            color_cls = COLORS[idx_it % len(COLORS)]
            cells += (
                f'<div class="vc">'
                f'<div class="ic">★</div>'
                f'<div class="bg {color_cls}">{esc(it["big"])}</div>'
                f'<div class="tt">{esc(it["title"])}</div>'
                f'<p>{it["desc"]}</p>'
                f'</div>'
            )
        return j - i, f'<div class="vg">{cells}</div>'
    return 0, ""


def try_ending(texts, i):
    """匹配结尾页：'从工具调用' + '到行为封装' + 副标题 + 多条要点"""
    if i >= len(texts) or texts[i] != "从工具调用":
        return 0, ""
    if i + 1 >= len(texts) or texts[i+1] != "到行为封装":
        return 0, ""
    j = i + 2
    subtitle = ""
    if j < len(texts):
        subtitle = texts[j]
        j += 1
    points = []
    while j < len(texts):
        t = texts[j]
        if is_watermark(t):
            j += 1
            continue
        if is_highlight(t):
            points.append({"type": "hl", "text": t})
        else:
            points.append({"type": "p", "text": t})
        j += 1
    html = (
        '<div class="end">'
        f'<h2>{esc(texts[i])} {esc(texts[i+1])}</h2>'
        f'<div class="sub" style="font-size:18px;color:var(--tm);margin-bottom:24px">{esc(subtitle)}</div>'
        '<div class="pts">'
    )
    for pt in points:
        if pt["type"] == "hl":
            html += f'<div class="hb g" style="margin:12px 0">{esc(pt["text"])}</div>'
        else:
            html += f'<p style="margin:8px 0">{esc(pt["text"])}</p>'
    html += '</div></div>'
    return j - i, html


def try_challenge_grid(texts, i):
    """匹配挑战卡片：[短标题 + 长描述] × N"""
    items = []
    j = i
    while j + 1 < len(texts):
        t1, t2 = texts[j], texts[j+1]
        if (is_short_title(t1) and not is_subtitle(t1) and not is_arch_layer(t1)
            and is_description(t2) and len(t1) <= 12 and len(t2) > 15
            and not is_code(t1) and not is_code(t2)
            and not is_stat_line(t1) and not is_combined_stat(t1)
            and not is_highlight(t2)):
            items.append({"title": t1, "desc": t2})
            j += 2
            continue
        break
    if len(items) >= 3:
        cells = ""
        for it in items:
            cells += (
                f'<div class="chc">'
                f'<h4>{esc(it["title"])}</h4>'
                f'<p>{esc(it["desc"])}</p>'
                f'</div>'
            )
        return j - i, f'<div class="chg">{cells}</div>'
    return 0, ""


def try_metrics_table(texts, i):
    """匹配量化对比表：方案/初始加载/单次请求/有效推理空间"""
    # 检测表头模式: "方案" + "初始加载" + "单次请求" + "有效推理空间"
    if not (i + 3 < len(texts)
            and texts[i] == "方案"
            and texts[i+1] == "初始加载"
            and texts[i+2] == "单次请求"
            and texts[i+3] == "有效推理空间"):
        return 0, ""
    j = i + 4
    rows = []
    # 全量加载行
    if j + 3 < len(texts) and texts[j] == "全量加载":
        rows.append({"name": "全量加载", "cells": [texts[j+1], texts[j+2], texts[j+3]], "type": "bad"})
        j += 4
    # 渐进式披露行
    if j + 3 < len(texts) and texts[j] == "渐进式披露":
        rows.append({"name": "渐进式披露", "cells": [texts[j+1], texts[j+2], texts[j+3]], "type": "good"})
        j += 4
    # 节省量行
    if j + 3 < len(texts) and texts[j] == "节省量":
        rows.append({"name": "节省量", "cells": [texts[j+1], texts[j+2], texts[j+3]], "type": "save"})
        j += 4
    if len(rows) >= 2:
        html = '<table class="tbl"><tr><th>方案</th><th>初始加载</th><th>单次请求</th><th>有效推理空间</th></tr>'
        for r in rows:
            style = ""
            if r["type"] == "good":
                style = ' style="color:var(--gl)"'
            elif r["type"] == "bad":
                style = ' style="color:var(--rd)"'
            elif r["type"] == "save":
                style = ' style="font-weight:600"'
            cells_html = "".join(f'<td{style}>{esc(c)}</td>' for c in r["cells"])
            html += f'<tr><td{style}>{esc(r["name"])}</td>{cells_html}</tr>'
        html += '</table>'
        return j - i, html
    return 0, ""


def try_benefit_row(texts, i):
    """匹配收益行：[短标签 + 描述] × 3
    如: 更低成本 | 节省 30-50% token 消耗
    标签是"更XX"或"X度"等2-4字短词，描述是带价值信息的句子"""
    items = []
    j = i
    # 跳过副标题
    start = i
    if start < len(texts) and is_subtitle(texts[start]) and ("收益" in texts[start] or "价值" in texts[start]):
        start += 1
    j = start
    while j + 1 < len(texts):
        t1, t2 = texts[j], texts[j+1]
        # t1 是短标签（≤6字，如"更低成本"/"更长推理"/"更清晰"），t2 是描述
        if (is_short_title(t1) and not is_subtitle(t1) and len(t1) <= 6
            and not is_arch_layer(t1) and not is_code(t1)
            and is_description(t2) and len(t2) > 6
            and not is_code(t2) and not is_highlight(t2)
            and not is_short_title(t2)):
            items.append({"label": t1, "desc": t2})
            j += 2
            continue
        break
    if len(items) >= 3:
        cells = ""
        for idx_it, it in enumerate(items):
            color_cls = COLORS[idx_it % len(COLORS)]
            cells += (
                f'<div class="st"><div class="sv {color_cls}" style="font-size:18px;font-weight:700">'
                f'{esc(it["label"])}</div><div class="sl">{esc(it["desc"])}</div></div>'
            )
        prefix = f'<h4>{esc(texts[i])}</h4>' if start > i else ""
        return j - i, prefix + f'<div class="sr">{cells}</div>'
    return 0, ""


def try_value_grid_b(texts, i):
    """匹配收益卡片：[标题 + 大字 + 描述] × N (无外层标题)
    标题是短标签（如"成本"、"质量"、"响应速"），大字是核心价值（如"节省 60-90%"）"""
    items = []
    j = i
    # 跳过可能的副标题（如"带来的三大收益"）
    start = i
    if start < len(texts) and is_subtitle(texts[start]) and "收益" in texts[start]:
        start += 1
    j = start
    # 三大收益模式: "成本" + "节省 60-90%" + "token 消耗..."
    while j + 2 < len(texts):
        t1, t2, t3 = texts[j], texts[j+1], texts[j+2]
        # t1 是短标签（≤5字），t2 是大字值（含数字/百分比/形容词），t3 是描述
        if (is_short_title(t1) and not is_subtitle(t1) and len(t1) <= 5
            and not is_arch_layer(t1) and not is_code(t1)
            and is_short_title(t2) and not is_subtitle(t2)
            and is_description(t3) and len(t3) > 8
            and not is_code(t3) and not is_highlight(t3)):
            items.append({"label": t1, "big": t2, "desc": t3})
            j += 3
            continue
        break
    if len(items) >= 2:
        cells = ""
        for idx_it, it in enumerate(items):
            color_cls = COLORS[idx_it % len(COLORS)]
            cells += (
                f'<div class="vc">'
                f'<div class="bg {color_cls}">{esc(it["big"])}</div>'
                f'<div class="tt">{esc(it["label"])}</div>'
                f'<p>{esc(it["desc"])}</p>'
                f'</div>'
            )
        # 如果跳过了副标题，把它作为 h4 前置
        prefix = f'<h4>{esc(texts[i])}</h4>' if start > i else ""
        return j - i, prefix + f'<div class="vg">{cells}</div>'
    return 0, ""


# ===== 主渲染逻辑 =====

def render_slide(slide):
    """渲染单页幻灯片 - 顺序模式匹配"""
    parts = []

    all_shapes = slide.get("shapes", [])
    text_shapes = [s for s in all_shapes if s.get("text", "").strip() and not is_watermark(s["text"])]
    table_shapes = [s for s in all_shapes if s.get("table")]

    if not text_shapes and not table_shapes:
        return ""

    texts = [s["text"].strip() for s in text_shapes]

    # 跳过 "Part X ·" 章节标识
    idx = 0
    if idx < len(texts) and re.match(r"Part\s*\d", texts[idx]):
        idx += 1

    # 结尾页特殊处理（"从工具调用" + "到行为封装"）
    # 检测模式：title="从工具调用" + remaining[0]="到行为封装"
    if idx + 1 < len(texts) and texts[idx] == "从工具调用" and texts[idx+1] == "到行为封装":
        end_texts = texts[idx:]
        consumed, html_str = try_ending(end_texts, 0)
        if consumed > 0:
            return html_str

    # 第一个剩余文本作为标题
    title = ""
    if idx < len(texts):
        title = texts[idx]
        idx += 1

    if title:
        parts.append(f"<h3>{esc(title)}</h3>")

    remaining = texts[idx:]

    if not remaining:
        for ts in table_shapes:
            parts.append(render_table(ts["table"]))
        return f'<div class="card">{"".join(parts)}</div>' if parts else ""

    # 顺序模式匹配
    i = 0
    while i < len(remaining):
        t = remaining[i]

        # 1. 代码块
        if is_code(t):
            parts.append(render_code(t))
            i += 1
            continue

        # 2. 高亮总结
        if is_highlight(t):
            parts.append(render_highlight(t))
            i += 1
            continue

        # 3. 量化对比表（"方案/初始加载/单次请求/有效推理空间" 表头）
        consumed, html_str = try_metrics_table(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 4. 策略对比卡片（"标题 + 优势 + 描述 + 不足 + 描述 + 适用 + 描述"）
        consumed, html_str = try_compare_grid(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 5. 价值卡片（标题 + 大字 + 标签 + 描述）
        consumed, html_str = try_value_grid(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 6. 收益卡片（标签 + 大字 + 描述）
        consumed, html_str = try_value_grid_b(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 6.5 收益行（短标签 + 带数字描述）
        consumed, html_str = try_benefit_row(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 7. 挑战卡片（短标题 + 长描述）
        consumed, html_str = try_challenge_grid(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 8. 步骤流程（"01\n标题" 格式，优先于统计行）
        consumed, html_str = try_step_flow(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 9. 概念表（短标题+描述对，优先于卡片网格）
        consumed, html_str = try_concept_table(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 10. 水平统计行（优先于垂直统计列表）
        consumed, html_str = try_stat_row(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 11. 卡片网格
        consumed, html_str = try_card_grid(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 12. 统计列表
        consumed, html_str = try_stat_list(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 13. 架构层级
        consumed, html_str = try_architecture(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 14. 表格
        consumed, html_str = try_table(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 15. 流程步骤
        consumed, html_str = try_flow_steps(remaining, i)
        if consumed > 0:
            parts.append(html_str)
            i += consumed
            continue

        # 16. 子标题
        if is_subtitle(t):
            parts.append(render_subtitle(t))
            i += 1
            continue

        # 17. 流程描述（含箭头）
        if has_arrow(t) and len(t) < 100:
            parts.append(render_flow([t]))
            i += 1
            continue

        # 18. 段落
        parts.append(f"<p>{esc(t)}</p>")
        i += 1

    # 渲染表格 shapes
    for ts in table_shapes:
        parts.append(render_table(ts["table"]))

    if not parts:
        return ""

    return f'<div class="card">{"".join(parts)}</div>'


def is_header(slide):
    texts = [s.get("text", "").strip() for s in slide.get("shapes", []) if s.get("text", "").strip()]
    if slide["index"] == 1:
        return True
    if texts and re.match(r"^\d{2}$", texts[0]):
        return True
    if "Part" in " ".join(texts) and len(texts) <= 6:
        return True
    return False


def get_title(slide):
    texts = [s.get("text", "").strip() for s in slide.get("shapes", [])
             if s.get("text", "").strip() and not is_watermark(s.get("text", ""))]
    if not texts:
        return f"第{slide['index']}页", ""

    if slide["index"] == 1:
        # 封面：找最长文本作为标题
        title = max(texts, key=len) if texts else "PPT 内容"
        subtitle = ""
        for t in texts:
            if t != title and len(t) > 5:
                subtitle = t
                break
        return title, subtitle

    pn, ti, su = "", "", ""
    for t in texts:
        if re.match(r"^\d{2}$", t):
            pn = t
        elif not ti and len(t) > 2:
            ti = t
        elif not su and len(t) > 2:
            su = t
    return ti or f"第{slide['index']}页", su


def build_html(data, name="PPT"):
    sections = []
    cur = None
    for slide in data:
        if is_header(slide):
            t, s = get_title(slide)
            cur = {"id": f"sec-{len(sections)}", "t": t, "s": s, "slides": []}
            sections.append(cur)
        else:
            if cur is None:
                cur = {"id": "sec-0", "t": "内容", "s": "", "slides": []}
                sections.append(cur)
            cur["slides"].append(slide)

    cover = sections[0] if sections else {"t": name, "s": ""}
    ct = cover["t"] or name
    cs = cover["s"]

    # 封面标签
    cover_tags = ""
    if sections:
        first_slide = data[0]
        tags = []
        for shape in first_slide.get("shapes", []):
            t = shape.get("text", "").strip()
            if t and not is_watermark(t) and t != ct and len(t) <= 15:
                tags.append(t)
        if tags:
            tag_classes = ["purple", "cyan", "green", "pink"]
            tag_html = []
            for i, t in enumerate(tags[:4]):
                cls = tag_classes[i % len(tag_classes)]
                tag_html.append(f'<span class="tag {cls}">{esc(t)}</span>')
            cover_tags = f'<div class="tags">{"".join(tag_html)}</div>'

    links = []
    side = []
    for i, sec in enumerate(sections):
        t = sec["t"][:12]
        links.append(f'<a href="#{sec["id"]}">{esc(t)}</a>')
        side.append(f'<div class="si" onclick="go(\'{sec["id"]}\')"><span class="n">{i:02d}</span> {esc(t)}</div>')

    content = [f'<section id="{sections[0]["id"]}" class="cover">{cover_tags}<h1>{esc(ct)}</h1><div class="sub">{esc(cs)}</div><div class="desc">共 {len(data)} 页幻灯片</div></section>']

    for i, sec in enumerate(sections[1:] if len(sections) > 1 else sections):
        parts = [f'<section id="{sec["id"]}"><div class="sh"><div class="pn">Section {i+1:02d}</div><h2>{esc(sec["t"])}</h2><div class="s">{esc(sec["s"])}</div></div>']
        for slide in sec["slides"]:
            r = render_slide(slide)
            if r:
                parts.append(r)
        parts.append("</section>")
        content.append("".join(parts))

    return (TEMPLATE.replace("__TITLE__", esc(ct))
            .replace("__BRAND__", esc(name))
            .replace("__LINKS__", "".join(links))
            .replace("__SIDE__", "".join(side))
            .replace("__CONTENT__", "".join(content)))


def main():
    ap = argparse.ArgumentParser(description="PPTX 转 HTML（高效智能版）")
    ap.add_argument("input", help="输入 PPTX 文件")
    ap.add_argument("-o", "--output", help="输出 HTML 路径")
    args = ap.parse_args()

    inp = Path(args.input)
    if not inp.exists():
        print(f"错误: 文件不存在: {inp}", file=sys.stderr)
        sys.exit(1)

    try:
        from pptx import Presentation
    except ImportError:
        print("错误: 需要安装 python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(inp))
    slides = []
    for i, slide in enumerate(prs.slides):
        shapes = []
        for shape in slide.shapes:
            info = {"type": str(shape.shape_type), "name": shape.name}
            if shape.has_text_frame:
                texts = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
                if texts:
                    info["text"] = "\n".join(texts)
            if shape.has_table:
                info["table"] = [[c.text for c in row.cells] for row in shape.table.rows]
            if info.get("text", "").strip() or info.get("table"):
                shapes.append(info)
        slides.append({"index": i + 1, "shapes": shapes})

    out = Path(args.output) if args.output else Path.cwd() / f"{inp.stem}.html"
    out.write_text(build_html(slides, inp.stem), encoding="utf-8")
    print(f"已生成: {out}")


if __name__ == "__main__":
    main()
