"""
读取 .pptx 文件文字内容（文本框 + 表格），逐页输出。
用法：python read_ppt.py "<ppt文件路径>"
"""
import sys
from pptx import Presentation


def read_ppt(filepath: str) -> str:
    """返回 PPT 全部文字内容，一次调用读取所有页面"""
    prs = Presentation(filepath)
    lines = []
    lines.append(f"文件：{filepath}")
    lines.append(f"共 {len(prs.slides)} 页\n")

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"{'='*60}")
        lines.append(f"  第 {i} 页 / 共 {len(prs.slides)} 页")
        lines.append(f"{'='*60}")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    lines.append(" | ".join(cells))
        lines.append("")

    return "\n".join(lines)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法：python read_ppt.py <ppt文件路径>")
        sys.exit(1)

    print(read_ppt(sys.argv[1]))
